#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
根据 architecture_diagram_analysis.md 的结构与颜色生成架构图 PPT。
布局：自上而下 5 个主层 + 右侧竖向「IDE和工具链」；每层按分析文档的行列与颜色绘制。

Usage:
    pip install python-pptx
    python3 build_architecture_ppt.py [output.pptx]
"""

import math
from pathlib import Path

# 从分析文档提取的层级结构与内容
LAYERS = [
    {
        "name": "应用层",
        "bg_rgb": (0xF0, 0xF0, 0xF0),   # 浅灰
        "border_rgb": (0xE0, 0xE0, 0xE0),
        "block_fill": (0xFF, 0xFF, 0xFF),
        "block_border": (0xE0, 0xE0, 0xE0),
        "rows": [
            ["系统应用", "桌面", "控制", "设置", "电话", "…"],
            ["OpenHarmony SDK", "Ability Kit", "ArkTS/合辑", "ArkData", "ArkUI", "ArkWeb…", "Xxx.kit"],
        ],
    },
    {
        "name": "应用框架层",
        "bg_rgb": (0xF0, 0xF0, 0xF0),
        "border_rgb": (0xE0, 0xE0, 0xE0),
        "block_fill": (0xFF, 0xFF, 0xFF),
        "block_border": (0xE0, 0xE0, 0xE0),
        "rows": [
            ["UI框架", "元能力框架"],
            ["用户程序框架"],
        ],
    },
    {
        "name": "系统服务层",
        "bg_rgb": (0xE6, 0xF0, 0xFA),   # 浅蓝
        "border_rgb": (0xE0, 0xE0, 0xE0),
        "block_fill": (0xFF, 0xFF, 0xFF),
        "block_border": (0xE0, 0xE0, 0xE0),
        "cols": [
            {
                "title": "系统基本能力子系统集",
                "left_col": ["分布式任务调度", "分布式数据管理", "分布式软总线", "方舟多语言运行时"],
                "right_col": ["多模输入子系统", "图形图像子系统", "安全子系统", "AI软件服务"],
                "bottom": ["公共基础类库子系统"],
            },
            {
                "title": "基础软件服务子系统集",
                "items": ["事件通知子系统", "通信子系统", "媒体子系统", "DFX子系统", "测试子系统", "…"],
            },
            {
                "title": "增强软件服务子系统集",
                "items": ["智慧屏业务子系统", "穿戴业务子系统", "IOT业务子系统", "…"],
            },
            {
                "title": "硬件服务子系统集",
                "items": ["位置服务子系统", "生物特征识别子系统", "穿戴专有硬件服务子系统", "IOT专有硬件服务子系统", "…"],
            },
        ],
    },
    {
        "name": "内核层",
        "bg_rgb": (0xE6, 0xF0, 0xFA),
        "border_rgb": (0xB8, 0xCC, 0xE4),
        "block_fill": (0xFF, 0xFF, 0xFF),
        "block_border": (0xB8, 0xCC, 0xE4),
        "rows": [
            ["KAL (内核抽象层)", "内核子系统", "Linux Kernel", "LiteOS", "Hongmeng"],
            ["驱动子系统", "HDF (统一驱动框架)"],
        ],
    },
    {
        "name": "芯片层",
        "bg_rgb": (0x6C, 0x6C, 0x6C),   # 深灰
        "border_rgb": (0xE0, 0xE0, 0xE0),
        "block_fill": (0xFF, 0xFF, 0xFF),
        "block_border": (0xE0, 0xE0, 0xE0),
        "rows": [
            ["海思平台", "L2 磐石开发板 (麒麟8000)", "L1 hispark_aiffy (HI3403V100)", "L0 恒玄、博流…"],
            ["其它芯片平台", "L2 展锐P7885、LoongArch…", "L1 RISC-V…", "L0 恒玄、博流…"],
        ],
    },
]

IDE_STRIP = {
    "name": "IDE和工具链",
    "bg_rgb": (0xE6, 0xF0, 0xFA),
    "text_rotation": -90,
}


def _rgb(rgb_tuple):
    from pptx.dml.color import RGBColor
    return RGBColor(*rgb_tuple)


def _add_rect(slide, left, top, width, height, fill_rgb, border_rgb, text="", font_pt=10, rotation=0):
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_rgb)
    if border_rgb:
        shape.line.color.rgb = _rgb(border_rgb)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(font_pt)
        p.font.bold = True
        p.font.name = "SimHei"
        p.font.color.rgb = _rgb((0, 0, 0))
        if rotation != 0:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def _render_simple_layer(slide, layer, left, top, width, height, pad=0.05, gap=0.03):
    """渲染应用层/应用框架层/内核层/芯片层：多行横向块。"""
    from pptx.util import Inches
    fill = layer["block_fill"]
    border = layer["block_border"]
    rows = layer.get("rows", [])
    if not rows:
        return
    n_rows = len(rows)
    row_h = (height - 2 * pad - (n_rows - 1) * gap) / n_rows
    y = top + pad
    for row_items in rows:
        n = len(row_items)
        if n == 0:
            y += row_h + gap
            continue
        cell_w = (width - 2 * pad - (n - 1) * gap) / n
        cell_w = max(cell_w, 0.15)
        x = left + pad
        for label in row_items:
            _add_rect(slide, x, y, cell_w, row_h, fill, border, label, font_pt=9)
            x += cell_w + gap
        y += row_h + gap


def _render_system_service_layer(slide, layer, left, top, width, height, pad=0.05, gap=0.03):
    """渲染系统服务层：四列，每列纵向或双列+底行。"""
    from pptx.util import Inches
    fill = layer["block_fill"]
    border = layer["block_border"]
    cols = layer.get("cols", [])
    if not cols:
        return
    n_cols = len(cols)
    col_w = (width - 2 * pad - (n_cols - 1) * gap) / n_cols
    x = left + pad
    for col in cols:
        if "items" in col:
            items = col["items"]
            n = len(items)
            cell_h = (height - 2 * pad - (n - 1) * gap) / n if n else 0.2
            cell_h = max(cell_h, 0.12)
            y = top + pad
            for label in items:
                _add_rect(slide, x, y, col_w, cell_h, fill, border, label, font_pt=8)
                y += cell_h + gap
        else:
            left_col = col.get("left_col", [])
            right_col = col.get("right_col", [])
            bottom = col.get("bottom", [])
            n_side = max(len(left_col), len(right_col), 1)
            cell_h = (height - 2 * pad - (n_side - 1) * gap - 0.2) / n_side if n_side else 0.15
            cell_h = max(cell_h, 0.1)
            half = (col_w - gap) / 2
            y = top + pad
            for i in range(max(len(left_col), len(right_col))):
                if i < len(left_col):
                    _add_rect(slide, x, y, half, cell_h, fill, border, left_col[i], font_pt=7)
                if i < len(right_col):
                    _add_rect(slide, x + half + gap, y, half, cell_h, fill, border, right_col[i], font_pt=7)
                y += cell_h + gap
            if bottom:
                bh = 0.18
                _add_rect(slide, x, y, col_w, bh, fill, border, bottom[0], font_pt=7)
        x += col_w + gap


def build_ppt(out_path: Path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = 13.333
    slide_h = 7.5
    margin = 0.2
    ide_width = 0.6
    content_w = slide_w - 2 * margin - ide_width - 0.05
    content_left = margin
    n_layers = len(LAYERS)
    layer_gap = 0.04
    total_gap = (n_layers - 1) * layer_gap
    layer_height = (slide_h - 2 * margin - total_gap) / n_layers

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 5 个主层
    for i, layer in enumerate(LAYERS):
        top = margin + i * (layer_height + layer_gap)
        # 层背景
        layer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(content_left), Inches(top), Inches(content_w), Inches(layer_height),
        )
        layer_shape.fill.solid()
        layer_shape.fill.fore_color.rgb = _rgb(layer["bg_rgb"])
        if layer.get("border_rgb"):
            layer_shape.line.color.rgb = _rgb(layer["border_rgb"])

        # 层标题（左侧窄条，无边框）
        title_w = 0.5
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(content_left), Inches(top), Inches(title_w), Inches(layer_height),
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = _rgb(layer["bg_rgb"])
        title_box.line.fill.background()
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        p.text = layer["name"]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.name = "SimHei"
        p.font.color.rgb = _rgb((0, 0, 0))

        inner_left = content_left + title_w + 0.03
        inner_w = content_w - title_w - 0.03
        inner_top = top
        inner_h = layer_height

        if "cols" in layer:
            _render_system_service_layer(slide, layer, inner_left, inner_top, inner_w, inner_h)
        else:
            _render_simple_layer(slide, layer, inner_left, inner_top, inner_w, inner_h)

    # 右侧 IDE和工具链（竖向贯穿）
    ide_left = content_left + content_w + 0.05
    ide_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(ide_left), Inches(margin), Inches(ide_width), Inches(slide_h - 2 * margin),
    )
    ide_shape.fill.solid()
    ide_shape.fill.fore_color.rgb = _rgb(IDE_STRIP["bg_rgb"])
    ide_shape.line.fill.background()
    tf = ide_shape.text_frame
    tf.word_wrap = True
    # 竖排：每字一行（逆时针约 90° 效果）
    vertical_text = "\n".join(IDE_STRIP["name"])
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = vertical_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.name = "SimHei"
    p.font.color.rgb = _rgb((0, 0, 0))

    prs.save(out_path)


def main():
    import sys
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "architecture.pptx"
    build_ppt(out)
    print("Written:", out, file=sys.stderr)


if __name__ == "__main__":
    main()
