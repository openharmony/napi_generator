#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
根据 MD 文档中的表格生成结构图 PPT：一行一个矩形框（深灰），一列为框内子模块（矩形），
支持「；」分多组、「 - 」分层级。文字左对齐、黑体、12 号字。

Usage:
    pip install python-pptx  # 依赖
    python3 ohppt.py <input.md> [output.pptx]
"""

import math
import re
import sys
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent

# 表格行：跳过表头与分隔行，解析 | 列1 | 列2 |
def _parse_md_table(md_path: Path) -> list[tuple[str, str]]:
    rows = []
    try:
        text = md_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return rows
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    for line in lines:
        if not line.startswith("|") or line.strip() == "|":
            continue
        parts = [p.strip() for p in line.split("|") if p.strip()]
        if len(parts) < 2:
            continue
        # 跳过表头行（列名含“行”“列”等）
        if "行" in parts[0] and "列" in parts[1]:
            continue
        if re.match(r"^[-:\s]+$", parts[0]) or re.match(r"^[-:\s]+$", parts[1]):
            continue
        rows.append((parts[0], parts[1]))
    return rows


def _parse_cell_content(cell: str) -> list[tuple[str, list[str]]]:
    """
    解析列内容：按「；」或「;」分多组；每组内按「 - 」拆分，第一个为子层名，后面均为该子层下的模块。
    例如「系统应用 - 示例 - 控制栏 - 设置 - 电话」→ 子层「系统应用」，模块 [示例, 控制栏, 设置, 电话]。
    返回 [(子层名, [模块1, 模块2, ...]), ...]
    """
    out = []
    for segment in re.split(r"[；;]", cell):
        segment = segment.strip()
        if not segment:
            continue
        parts = [p.strip() for p in re.split(r"\s+-\s+", segment) if p.strip()]
        if not parts:
            continue
        sub_layer_name = parts[0]
        modules = parts[1:] if len(parts) > 1 else []
        out.append((sub_layer_name, modules))
    return out


def _parse_row_label(row_label: str) -> tuple[bool, str | None, str]:
    """
    解析行标签，识别父子层级。
    - 「8. 内核层（标题）」→ 层行，返回 (False, None, "内核层")
    - 「9. 内核层-内嵌子系统」→ 子层行，返回 (True, "内核层", "内嵌子系统")
    返回 (is_sub_layer, parent_name_or_None, layer_or_sub_name)
    """
    main = re.sub(r"^\d+\.\s*", "", row_label).strip()
    if "（标题）" in main:
        name = main.replace("（标题）", "").strip()
        return (False, None, name)
    if "-" in main:
        parts = main.split("-", 1)
        left, right = parts[0].strip(), parts[1].strip()
        if left.endswith("层"):
            return (True, left, right)
    return (False, None, main)


def _build_layer_blocks(rows: list[tuple[str, str]]) -> list[dict]:
    """
    按行标签的父子关系合并行，得到层块列表。
    每个层块: {"name": 匹配用层名, "display": 显示用标题, "sub_layers": [(子层名, [模块词...]), ...]}
    """
    blocks: list[dict] = []
    for row_label, content in rows:
        is_sub, parent_name, self_name = _parse_row_label(row_label)
        if is_sub:
            groups = _parse_cell_content(content)
            items = (groups[0][1] if groups else [])
            if not items and groups:
                items = [groups[0][0]]
            for b in reversed(blocks):
                if b["name"] == parent_name:
                    b["sub_layers"].append((self_name, items))
                    break
            else:
                blocks.append({
                    "name": parent_name,
                    "display": parent_name,
                    "sub_layers": [(self_name, items)],
                })
        else:
            groups = _parse_cell_content(content)
            if len(groups) == 1 and groups[0][0] == self_name and len(groups[0][1]) == 0:
                groups = []
            sub_layers = [(g[0], g[1]) for g in groups]
            display = re.sub(r"^\d+\.\s*", "", row_label).strip()
            if display.endswith("（标题）"):
                display = display.replace("（标题）", "").strip()
            blocks.append({
                "name": self_name,
                "display": display[:80],
                "sub_layers": sub_layers,
            })
    return blocks


# 按字数估算框宽（英寸）
INCH_PER_CHAR = 0.058
# 上下左右留 10 的间隔（0.1 英寸）
PAD = 0.05
# 每个框最小宽度、最小高度（单位：pt），1 inch = 72 pt
MIN_WIDTH_PT = 20
MIN_HEIGHT_PT = 10
MIN_WIDTH_IN = MIN_WIDTH_PT / 72.0
MIN_HEIGHT_IN = MIN_HEIGHT_PT / 72.0
# 一行最多放 6 个子模块，不够的换行；模块宽度为子层模块区宽度平均分配（按 6 列）
MAX_MODULES_PER_ROW = 6


def _module_layout_max6(
    sub_items: list[str],
    area_left: float,
    area_top: float,
    area_w: float,
    cell_h: float,
    mod_gap: float,
) -> tuple[list[tuple[float, float, float, float, str]], float]:
    """
    子层内：模块宽度 = 子层方框（模块区）宽度按 6 列平均分配，一行最多 6 个，不够换行。
    返回 ([(left, top, width, height, word), ...], 模块区所需高度)
    """
    if not sub_items:
        return [], 0.0
    n = len(sub_items)
    n_rows = math.ceil(n / MAX_MODULES_PER_ROW)
    cell_w = (area_w - (MAX_MODULES_PER_ROW - 1) * mod_gap) / MAX_MODULES_PER_ROW
    cell_w = max(cell_w, MIN_WIDTH_IN)
    required_h = n_rows * cell_h + (n_rows - 1) * mod_gap if n_rows else 0.0
    out: list[tuple[float, float, float, float, str]] = []
    for i, word in enumerate(sub_items):
        row_i, col_i = i // MAX_MODULES_PER_ROW, i % MAX_MODULES_PER_ROW
        left = area_left + col_i * (cell_w + mod_gap)
        top = area_top + row_i * (cell_h + mod_gap)
        out.append((left, top, cell_w, cell_h, word))
    return out, required_h


def _set_para_black(p, text: str, font_pt: int = 12, bold: bool = True):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(font_pt)
    p.font.bold = bold
    p.font.name = "SimHei"
    p.font.color.rgb = RGBColor(0, 0, 0)


def _create_ppt_from_table(layer_blocks: list[dict], out_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w_in = 13.333
    slide_h_in = 7.5

    BLUE = RGBColor(0x00, 0x4A, 0xAD)
    LIGHT_BLUE = RGBColor(0xE6, 0xF0, 0xFA)
    # 方框边框颜色淡于填充色
    LIGHT_BORDER = RGBColor(0xB8, 0xCC, 0xE4)
    # 模块：淡灰填充 + 淡灰边框
    MODULE_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
    MARGIN = 0.2
    LAYER_GAP = 0.05
    INNER_MARGIN = 0
    SUB_LAYER_GAP = 0.02
    LAYER_TITLE_W = 1.0
    SUB_TITLE_W = 0.9
    TITLE_H = 0.18
    MODULE_GAP = 0.03
    FONT_PT = 10

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    n_blocks = len(layer_blocks)
    avail_h = slide_h_in - 2 * MARGIN - (n_blocks - 1) * LAYER_GAP
    layer_height = avail_h / n_blocks if n_blocks else avail_h

    for row_idx, block in enumerate(layer_blocks):
        layer_top = MARGIN + row_idx * (layer_height + LAYER_GAP)
        row_left = Inches(MARGIN)
        row_top = Inches(layer_top)
        row_width = Inches(slide_w_in - 2 * MARGIN)
        row_height = Inches(layer_height)

        layer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            row_left, row_top, row_width, row_height
        )
        layer_shape.fill.solid()
        layer_shape.fill.fore_color.rgb = LIGHT_BLUE
        layer_shape.line.color.rgb = LIGHT_BORDER

        display = block["display"]
        sub_layers = block["sub_layers"]
        # 层内上下左右留 PAD 间隔，内容尽量平铺在父层方框内
        in_left = MARGIN + INNER_MARGIN + PAD
        in_top = layer_top + INNER_MARGIN + PAD
        in_w = slide_w_in - 2 * (MARGIN + INNER_MARGIN) - 2 * PAD
        in_h = layer_height - 2 * INNER_MARGIN - 2 * PAD
        cur_y = in_top

        layer_title_w = max(LAYER_TITLE_W, min(len(display) * INCH_PER_CHAR, 2.0))
        layer_title_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(in_left), Inches(cur_y), Inches(layer_title_w), Inches(in_h)
        )
        layer_title_box.fill.solid()
        layer_title_box.fill.fore_color.rgb = LIGHT_BLUE
        layer_title_box.line.fill.background()
        lt_tf = layer_title_box.text_frame
        lt_tf.word_wrap = True
        lp0 = lt_tf.paragraphs[0] if lt_tf.paragraphs else lt_tf.add_paragraph()
        _set_para_black(lp0, display[:120] if len(display) > 120 else display, font_pt=FONT_PT)

        content_left = in_left + layer_title_w + 0.03
        content_w = in_w - (layer_title_w + 0.03)

        if not sub_layers:
            no_gr_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(content_left), Inches(cur_y),
                Inches(max(content_w, MIN_WIDTH_IN)), Inches(max(min(0.4, in_h - 0.05), MIN_HEIGHT_IN))
            )
            no_gr_box.fill.solid()
            no_gr_box.fill.fore_color.rgb = LIGHT_BLUE
            no_gr_box.line.color.rgb = LIGHT_BORDER
            tf = no_gr_box.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            _set_para_black(p0, block.get("name", "")[:400], font_pt=FONT_PT)
            continue

        # 子层高度由模块行数决定：一行最多 6 个模块，行数越多子层越高
        K = len(sub_layers)
        n_rows_per_sub: list[int] = []
        for _g, sub_items in sub_layers:
            n = len(sub_items)
            n_rows_per_sub.append(math.ceil(n / MAX_MODULES_PER_ROW) if n else 1)
        total_rows = sum(n_rows_per_sub)
        total_gap = sum((r - 1) * MODULE_GAP for r in n_rows_per_sub)
        available_h = in_h - 2 * K * PAD - (K - 1) * SUB_LAYER_GAP - total_gap
        cell_h = (available_h / total_rows) if total_rows else MIN_HEIGHT_IN
        cell_h = max(cell_h, MIN_HEIGHT_IN)

        for idx, (group_title, sub_items) in enumerate(sub_layers):
            n_rows = n_rows_per_sub[idx]
            sub_layer_h = 2 * PAD + n_rows * cell_h + (n_rows - 1) * MODULE_GAP
            sub_layer_w = content_w

            sub_layer_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(content_left), Inches(cur_y), Inches(sub_layer_w), Inches(sub_layer_h)
            )
            sub_layer_box.fill.solid()
            sub_layer_box.fill.fore_color.rgb = LIGHT_BLUE
            sub_layer_box.line.color.rgb = LIGHT_BORDER

            sub_title_w = max(SUB_TITLE_W, min(len(group_title) * INCH_PER_CHAR, 1.8))
            mod_area_left = content_left + PAD + sub_title_w + 0.02
            mod_area_w_use = sub_layer_w - 2 * PAD - sub_title_w - 0.02

            sub_title_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(content_left + PAD), Inches(cur_y + PAD),
                Inches(sub_title_w), Inches(sub_layer_h - 2 * PAD)
            )
            sub_title_shape.fill.solid()
            sub_title_shape.fill.fore_color.rgb = LIGHT_BLUE
            sub_title_shape.line.color.rgb = LIGHT_BORDER
            tf0 = sub_title_shape.text_frame
            tf0.word_wrap = True
            p0 = tf0.paragraphs[0] if tf0.paragraphs else tf0.add_paragraph()
            _set_para_black(p0, group_title, font_pt=FONT_PT)

            mod_boxes, _ = _module_layout_max6(
                sub_items, mod_area_left, PAD, mod_area_w_use, cell_h, MODULE_GAP,
            )
            for left_i, top_i, w_i, h_i, word in mod_boxes:
                box = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(left_i), Inches(cur_y + top_i), Inches(w_i), Inches(h_i)
                )
                box.fill.solid()
                box.fill.fore_color.rgb = MODULE_GRAY
                box.line.color.rgb = MODULE_GRAY
                tb = box.text_frame
                tb.word_wrap = True
                tp = tb.paragraphs[0] if tb.paragraphs else tb.add_paragraph()
                _set_para_black(tp, word, font_pt=FONT_PT)
            cur_y += sub_layer_h + SUB_LAYER_GAP

    prs.save(out_path)


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: python3 ohppt.py <input.md> [output.pptx]", file=sys.stderr)
        return 1
    md_path = Path(sys.argv[1])
    if not md_path.is_file():
        print("Not a file:", md_path, file=sys.stderr)
        return 1
    out_path = Path(sys.argv[2]) if len(sys.argv) > 2 else md_path.with_suffix(".pptx")

    rows = _parse_md_table(md_path)
    if not rows:
        print("No table rows found in", md_path, file=sys.stderr)
        return 1
    layer_blocks = _build_layer_blocks(rows)
    try:
        _create_ppt_from_table(layer_blocks, out_path)
    except Exception as e:
        print("Error generating PPT:", e, file=sys.stderr)
        try:
            import traceback
            traceback.print_exc()
        except Exception:
            pass
        return 1
    print("Written:", out_path, "({} layers from {} rows)".format(len(layer_blocks), len(rows)), file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
